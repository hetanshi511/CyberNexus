import os
import logging
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_DIR = os.path.join(BASE_DIR, "templates")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# --- HELPER: MANUALLY STYLE TABLE ---
def style_table(table):
    """
    Applies a professional corporate look to the table.
    Handles errors gracefully.
    """
    try:
        if not table or not table.rows:
            logger.warning("style_table: Invalid table object")
            return

        # 1. Style the Header Row (Row 0)
        if len(table.rows) == 0:
            logger.warning("style_table: Table has no rows")
            return

        header_row = table.rows[0]
        for cell in header_row.cells:
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Corporate Navy Blue

                if cell.text_frame.paragraphs:
                    p = cell.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER

            except Exception as e:
                logger.warning(f"style_table: Error styling header cell: {e}")
                continue

        # 2. Style Body Rows (Rows 1+)
        for i, row in enumerate(table.rows):
            if i == 0:
                continue  # Skip the header row we just styled

            for cell in row.cells:
                try:
                    if cell.text_frame.paragraphs:
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(14)
                        p.alignment = PP_ALIGN.LEFT

                except Exception as e:
                    logger.warning(f"style_table: Error styling body cell in row {i}: {e}")
                    continue

    except Exception as e:
        logger.error(f"style_table: Unexpected error: {e}", exc_info=True)



# --- HELPER: CHART CREATOR ---
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

# --- HELPER: CHART CREATOR ---
def create_chart(slide, chart_type, chart_data, description=None):
    """
    Creates a chart (Bar, Pie, or Line) on the slide based on chart_type and chart_data.
    Optionally adds a description below the chart.
    Returns True if successful, False otherwise.
    """
    if not slide:
        logger.error("create_chart: Invalid slide object")
        return False
    if not chart_data or not isinstance(chart_data, dict):
        logger.error("create_chart: Invalid chart_data")
        return False

    categories = chart_data.get('categories', [])
    values = chart_data.get('values', [])
    percentages = chart_data.get('percentages', [])
    series_name = chart_data.get('series_name', 'Series 1')

    if not categories or not values:
        logger.warning("create_chart: Missing categories or values")
        return False

    # --- FIX: DATA SANITIZATION ---
    # LLMs often send numbers as strings (e.g., "100"). This makes Line Charts blank.
    # We try to convert them to numbers before plotting.
    try:
        cleaned_values = []
        for v in values:
            if isinstance(v, (int, float)):
                cleaned_values.append(v)
            elif isinstance(v, str):
                # Remove currency symbols or commas if present
                clean_v = v.replace(',', '').replace('$', '').replace('%', '').strip()
                if '.' in clean_v:
                    cleaned_values.append(float(clean_v))
                else:
                    cleaned_values.append(int(clean_v))
            else:
                cleaned_values.append(0)
        values = cleaned_values
    except Exception as e:
        logger.warning(f"create_chart: Data sanitization failed, using original values. Error: {e}")

    # Align lengths
    if len(categories) != len(values):
        logger.warning(f"create_chart: Mismatch - {len(categories)} categories vs {len(values)} values")
        min_len = min(len(categories), len(values))
        categories = categories[:min_len]
        values = values[:min_len]

    try:
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = categories
        chart_data_obj.add_series(series_name, values)
        
        # Define Position (Reduced height to make room for description)
        x, y, cx, cy = Inches(0.8), Inches(2.0), Inches(7.0), Inches(4.2)
        
        # Add Description if provided
        if description:
            try:
                desc_left = Inches(1.0)
                desc_top = Inches(6.3) # Below the chart (2.0 + 4.2 = 6.2)
                desc_width = Inches(8.0)
                desc_height = Inches(1.2) # Increased height
                
                txBox = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = description
                p.font.size = Pt(11) # Slightly smaller for more text
                # p.font.italic = True # Removed italic for readability
                p.alignment = PP_ALIGN.LEFT # Left align for paragraphs
            except Exception as e:
                logger.warning(f"create_chart: Failed to add description: {e}")
        
        chart_shape = None

        # --- CHART TYPE SELECTION ---
        if chart_type == "BAR_CHART":
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_obj
            )
        elif chart_type == "PIE_CHART":
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data_obj
            )
        elif chart_type == "LINE_CHART":
            # Use LINE_MARKERS so data points are visible even if the line is faint
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data_obj
            )
        else:
            logger.warning(f"create_chart: Unknown chart_type: {chart_type}")
            return False

        # --- POST-PROCESSING STYLING ---
        if chart_shape and chart_shape.chart:
            chart = chart_shape.chart
            
            # 1. Legend Configuration
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False # Check: prevents shrinking plot area too much
            
            try:
                chart.legend.font.size = Pt(12)
            except Exception as e:
                logger.warning(f"create_chart: Error adjusting legend font: {e}")

            # 2. Axis Titles (For Bar and Line Charts)
            if chart_type in ["BAR_CHART", "LINE_CHART"]:
                try:
                    # Y-axis (Value axis)
                    if chart.value_axis:
                        chart.value_axis.has_title = True
                        chart.value_axis.axis_title.text_frame.text = str(series_name)
                        if chart.value_axis.axis_title.text_frame.paragraphs:
                            p = chart.value_axis.axis_title.text_frame.paragraphs[0]
                            p.font.size = Pt(11)
                            p.font.bold = True

                    # X-axis (Category axis)
                    if chart.category_axis:
                        chart.category_axis.has_title = True
                        # If the agent provided a specific X-axis label in data, we could use it, 
                        # otherwise default to 'Categories' or leave blank to save space
                        chart.category_axis.axis_title.text_frame.text = "Categories" 
                        if chart.category_axis.axis_title.text_frame.paragraphs:
                            p = chart.category_axis.axis_title.text_frame.paragraphs[0]
                            p.font.size = Pt(11)
                except Exception as e:
                    logger.warning(f"create_chart: Error setting axis titles: {e}")

            # 3. Data Labels
            if chart.plots:
                plot = chart.plots[0]
                plot.has_data_labels = True
                
                # For Pie Charts: Show only percentages
                if chart_type == "PIE_CHART":
                    try:
                        data_labels = plot.data_labels
                        data_labels.show_percentage = True
                        data_labels.number_format = "0.0%"
                        # Hide values for pie charts, show only percentages
                        data_labels.show_value = False
                    except Exception as e:
                        logger.warning(f"create_chart: Error setting pie labels: {e}")
                
                # For Bar/Line Charts: Values
                else:
                    try:
                        data_labels = plot.data_labels
                        data_labels.number_format = "#,##0"
                        if hasattr(data_labels, 'font'):
                            data_labels.font.size = Pt(9)
                    except Exception as e:
                        logger.warning(f"create_chart: Error formatting data labels: {e}")

        logger.info(f"create_chart: Successfully created {chart_type}")
        return True

    except Exception as e:
        logger.error(f"create_chart: Error creating chart: {e}", exc_info=True)
        return False




# --- HELPER: TABLE CREATOR ---

def create_table(slide, table_data, description=None):

    """

    Creates a table on the slide based on table_data.

    Returns True if successful, False otherwise.

    """

    if not slide:

        logger.error("create_table: Invalid slide object")

        return False

   

    if not table_data or not isinstance(table_data, dict):

        logger.error("create_table: Invalid table_data")

        return False

   

    columns = table_data.get('columns', [])

    rows = table_data.get('rows', [])

   

    if not columns:

        logger.warning("create_table: No columns provided")

        return False

   

    if not rows:

        logger.warning("create_table: No rows provided")

        return False

   

    try:

        rows_count = len(rows) + 1  # +1 for header

        cols_count = len(columns)

       

        # Validate dimensions

        if rows_count < 1 or cols_count < 1:

            logger.error("create_table: Invalid table dimensions")

            return False

       

        # Position with improved margins to prevent overflow

        left = Inches(1.2)  # Increased from 0.8 to move further from left side

        top = Inches(2.0)

        width = Inches(8.0)  # Reduced slightly to maintain right margin

        height = Inches(0.8 * rows_count)  # Dynamic height based on rows

       

        shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)

        table = shape.table

       

        # 1. Populate Header

        for i, col_name in enumerate(columns):

            if i < cols_count:

                try:

                    cell = table.cell(0, i)

                    cell.text = str(col_name)

                except Exception as e:

                    logger.warning(f"create_table: Error populating header cell {i}: {e}")

                    continue

       

        # 2. Populate Rows

        for row_idx, row_data in enumerate(rows):

            if row_idx + 1 >= rows_count:

                break

            for col_idx, cell_data in enumerate(row_data):

                if col_idx < cols_count:

                    try:

                        cell = table.cell(row_idx + 1, col_idx)

                        cell.text = str(cell_data)

                    except Exception as e:

                        logger.warning(f"create_table: Error populating cell ({row_idx+1}, {col_idx}): {e}")

                        continue

       

        # 3. Apply Styling

        style_table(table)

       

        # 4. Add text wrapping to prevent overflow

        try:

            for row in table.rows:

                for cell in row.cells:

                    if cell.text_frame:

                        cell.text_frame.word_wrap = True

                        # Add small margins for better spacing

                        try:

                            cell.margin_left = Inches(0.02)

                            cell.margin_right = Inches(0.02)

                        except:

                            pass  # Ignore if margins can't be set

        except Exception as e:

            logger.warning(f"create_table: Error setting text wrapping: {e}")

       

        # 5. Add Description (if provided)
        if description:
            try:
                # Calculate position: below table or fixed minimum
                table_bottom_inches = 2.0 + (0.8 * rows_count) # top + height
                desc_top = max(table_bottom_inches + 0.2, 6.5)
                
                # Check if it goes off slide (approx 7.5 inches)
                if desc_top > 7.0:
                     logger.warning("create_table: Table too tall, skipping description to avoid overflow.")
                else:
                    desc_left = Inches(1.0)
                    desc_width = Inches(8.0)
                    desc_height = Inches(1.2)

                    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(desc_top), desc_width, desc_height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = description
                    p.font.size = Pt(11)
                    # p.font.italic = True
                    p.alignment = PP_ALIGN.LEFT
            except Exception as e:
                logger.warning(f"create_table: Failed to add description: {e}")

        logger.info(f"create_table: Successfully created table with {rows_count} rows and {cols_count} columns")

        return True

       

    except Exception as e:

        logger.error(f"create_table: Error creating table: {e}", exc_info=True)

        return False






# --- HELPER: ADD BRANDING (LOGO & PAGE NUMBERS) ---
def add_branding_to_slide(slide, slide_number, total_slides):
    """
    Adds company logo (Top-Left) and Page Number (Bottom-Right) to a slide.
    """
    try:
        # 1. Add Logo
        logo_path = os.path.join(TEMPLATE_DIR, "logo.png")
        if os.path.exists(logo_path):
            left = Inches(0.2)
            top = Inches(0.2)
            width = Inches(2.04) # Increased size
            # Check if logo already exists? Hard to check, just add.
            slide.shapes.add_picture(logo_path, left, top, width=width)
            
        # 2. Add Page Number
        # Position: Bottom Right
        # SKIP if it is the first slide (Title Slide)
        if slide_number > 1:
            left_tx = Inches(9.0)
            top_tx = Inches(7.1) 
            width_tx = Inches(0.8)
            height_tx = Inches(0.4)
            
            txBox = slide.shapes.add_textbox(left_tx, top_tx, width_tx, height_tx)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"{slide_number}"
            p.font.size = Pt(18) # Increased font size
            p.alignment = PP_ALIGN.RIGHT
        
    except Exception as e:
        logger.warning(f"add_branding_to_slide: Error adding branding: {e}")

# --- HELPER: DOCUMENT CONTROL SLIDE ---
def create_document_control_slide(prs, metadata):
    """
    Creates the Document Control slide at the specified index (usually 1).
    """
    try:
        layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(layout)
        
        # Title
        if slide.shapes.title:
            slide.shapes.title.text = "Document Control"
            
        # Create Table
        rows = 8
        cols = 2
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(7.0)
        height = Inches(0.5 * rows)
        
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        
        # Define fields
        fields = [
            ("Document Title", metadata.get("title")),
            ("Version", metadata.get("version")),
            ("Creation Date", metadata.get("creation_date")),
            ("Document Owner", metadata.get("owner")),
            ("Email-Contact", metadata.get("email")),
            ("Department", metadata.get("department")),
            ("Approved By", metadata.get("approved_by")),
            ("Classification", metadata.get("classification", "Internal"))
        ]
        
        for i, (label, value) in enumerate(fields):
            # Label Cell
            cell_label = table.cell(i, 0)
            cell_label.text = label
            cell_label.fill.solid()
            cell_label.fill.fore_color.rgb = RGBColor(220, 230, 240) # Light color
            if cell_label.text_frame.paragraphs:
                p = cell_label.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0) # Black text
                
            # Value Cell
            cell_value = table.cell(i, 1)
            cell_value.text = str(value)
            if cell_value.text_frame.paragraphs:
                p = cell_value.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0)
                
        # Move slide to index 1?
        # python-pptx appends slides. If we call this immediately after Slide 0, it will be Slide 1.
        return slide
    except Exception as e:
        logger.error(f"create_document_control_slide: Error: {e}")
        return None

# --- AGENT 1: VISUALIZATION BUILDER NODE ---

def visualization_builder_node(state):

    """

    Visualization Builder Agent: Creates charts and tables in PowerPoint.

   

    RESPONSIBILITIES:

    - Initialize PowerPoint presentation

    - Create title slide

    - Create slides with charts (BAR_CHART, PIE_CHART)

    - Create slides with tables

    - Save to temporary file for content builder

   

    INSTRUCTIONS:

    1. Load or create presentation template

    2. Create title slide (first slide) with title and subtitle

    3. For each slide with chart_type and chart_data:

       - Use layout_index 5 (Title Only)

       - Create appropriate chart type (BAR_CHART or PIE_CHART)

       - Position chart properly (avoid title overlap)

       - Add legend and data labels

    4. For each slide with table_data:

       - Use layout_index 5 (Title Only)

       - Create table with proper dimensions

       - Style header row (navy blue background, white text)

       - Style body rows (readable font size)

    5. Skip slides that only have text content (content builder will handle)

    6. Save to temporary PPTX file

    7. Return temp file path for content builder

    """

    logger.info("=== VISUALIZATION BUILDER AGENT: Creating charts and tables ===")

   

    try:

        topic = state.get("topic", "Presentation")

        slides_data = state.get("slides", [])

       

        # Validation

        if not slides_data:

            error_msg = "No slides data provided to visualization builder"

            logger.error(error_msg)

            return {

                "temp_pptx_path": "",

                "current_step": "visualization_build_failed",

                "error_message": error_msg

            }

       

        if len(slides_data) < 4:

            error_msg = f"Only {len(slides_data)} slides provided. Minimum 4 required."

            logger.error(error_msg)

            return {

                "temp_pptx_path": "",

                "current_step": "visualization_build_failed",

                "error_message": error_msg

            }

       

        logger.info(f"Visualization builder: Processing {len(slides_data)} slides")


        # Load template or create new presentation
        template_name = state.get("template_name", "corporate_blue")
        if not template_name.endswith(".pptx"):
            template_name += ".pptx"
            
        template_path = os.path.join(TEMPLATE_DIR, template_name)

        
        # Fallback if specific file doesn't exist
        if not os.path.exists(template_path):
             logger.warning(f"Template {template_name} not found. Falling back to light_red.pptx")
             # Load template or create new presentation
        # Load template or create new presentation
        template_name = state.get("template_name", "light_red")
        if not template_name.endswith(".pptx"):
            template_name += ".pptx"
            
        template_path = os.path.join(TEMPLATE_DIR, template_name)

        
        # Fallback if specific file doesn't exist
        if not os.path.exists(template_path):
             logger.warning(f"Template {template_name} not found. Falling back to light_blue.pptx")
             template_path = os.path.join("templates", "light_blue.pptx")

        prs = None

       

        if os.path.exists(template_path):

            try:

                prs = Presentation(template_path)

                logger.info(f"Loaded template from {template_path}")

            except Exception as e:

                logger.warning(f"Failed to load template: {e}. Using default template.")

                prs = Presentation()

        else:

            logger.info("Template not found. Using default template.")

            prs = Presentation()

       

        if not prs:

            error_msg = "Failed to initialize presentation"

            logger.error(error_msg)

            return {

                "temp_pptx_path": "",

                "current_step": "visualization_build_failed",

                "error_message": error_msg

            }

       

        # Track slides created

        slides_created = 0

        errors = []

        slide_mapping = {}  # Map slide index to actual slide object for content builder

       

        try:

            for index, slide_data in enumerate(slides_data):

                try:

                    if not isinstance(slide_data, dict):

                        errors.append(f"Slide {index + 1}: Invalid slide data structure")

                        continue

                   

                    title = slide_data.get('title', f'Slide {index + 1}')

                    is_chart = 'chart_type' in slide_data and slide_data.get('chart_type')

                    is_table = 'table_data' in slide_data and slide_data.get('table_data')

                    has_content = 'content' in slide_data and slide_data.get('content')

                   

                    # --- A. TITLE SLIDE (Slide 0) ---

                    if index == 0:

                        try:

                            layout_idx = 0

                            if layout_idx >= len(prs.slide_layouts):

                                layout_idx = len(prs.slide_layouts) - 1

                           

                            layout = prs.slide_layouts[layout_idx]

                            slide = prs.slides.add_slide(layout)

                            slide_mapping[index] = slide

                           

                            # Title

                            if slide.shapes.title:

                                slide.shapes.title.text = str(title)

                                # Set title font size to 48-54 points for first slide

                                try:

                                    for paragraph in slide.shapes.title.text_frame.paragraphs:

                                        paragraph.font.size = Pt(50)  # Midpoint of 48-54 range

                                        paragraph.font.bold = True

                                except Exception as e:

                                    logger.warning(f"Error setting title font size for first slide: {e}")

                           

                            # Subtitle (will be filled by content builder if needed)

                            if len(slide.placeholders) > 1:

                                try:

                                    subtitle = slide.placeholders[1]

                                    content = slide_data.get('content', [])

                                    if content and isinstance(content, list) and len(content) > 0:

                                        subtitle.text = str(content[0])

                                    elif content and isinstance(content, str):

                                        subtitle.text = content

                                    else:

                                        subtitle.text = ""

                                except Exception as e:

                                    logger.warning(f"Error setting subtitle for slide {index + 1}: {e}")

                           

                            slides_created += 1

                            logger.info(f"Created title slide: {title}")

                            # --- INSERT DOCUMENT CONTROL SLIDE (Slide 2) ---
                            # Only do this ONCE, immediately after title slide (index 0)
                            try:
                                meta = state.get("document_metadata", {})
                                # Set title in metadata if not present
                                if not meta.get("title"):
                                    meta["title"] = title
                                    
                                dc_slide = create_document_control_slide(prs, meta)
                                if dc_slide:
                                    logger.info("Created Document Control slide")
                                    slides_created += 1
                                    # Note: We do NOT add this to slide_mapping because there is no corresponding 'slide_data' index for it.
                                    # Content builder loop uses indices from slides_data.
                            except Exception as e:
                                logger.warning(f"Failed to create document control slide: {e}")

                            


                           

                        except Exception as e:

                            error_msg = f"Error creating title slide: {e}"

                            logger.error(error_msg)

                            errors.append(error_msg)

                            continue

                   

                    # --- B. CHART SLIDES ---

                    elif is_chart:

                        try:

                            layout_idx = 5  # Title Only

                            if layout_idx >= len(prs.slide_layouts):

                                layout_idx = min(5, len(prs.slide_layouts) - 1)

                           

                            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

                            slide_mapping[index] = slide

                           

                            # Set title

                            if slide.shapes.title:

                                slide.shapes.title.text = str(title)

                           

                            # Create chart

                            chart_type = slide_data.get('chart_type')

                            chart_data = slide_data.get('chart_data', {})
                            
                            description = slide_data.get('description')

                            success = create_chart(slide, chart_type, chart_data, description)

                           

                            if success:

                                slides_created += 1

                                logger.info(f"Created chart slide {index + 1}: {title} ({chart_type})")

                            else:

                                errors.append(f"Slide {index + 1} ('{title}'): Chart creation failed")

                               

                        except Exception as e:

                            error_msg = f"Error creating chart slide {index + 1}: {e}"

                            logger.error(error_msg, exc_info=True)

                            errors.append(error_msg)

                            continue

                   

                    # --- C. TABLE SLIDES ---

                    elif is_table:

                        try:

                            layout_idx = 5  # Title Only

                            if layout_idx >= len(prs.slide_layouts):

                                layout_idx = min(5, len(prs.slide_layouts) - 1)

                           

                            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

                            slide_mapping[index] = slide

                           

                            # Set title

                            if slide.shapes.title:

                                slide.shapes.title.text = str(title)

                           

                            # Create table

                            table_data = slide_data.get('table_data', {})
                            
                            description = slide_data.get('description')

                            success = create_table(slide, table_data, description)

                           

                            if success:

                                slides_created += 1

                                logger.info(f"Created table slide {index + 1}: {title}")

                            else:

                                errors.append(f"Slide {index + 1} ('{title}'): Table creation failed")

                               

                        except Exception as e:

                            error_msg = f"Error creating table slide {index + 1}: {e}"

                            logger.error(error_msg, exc_info=True)

                            errors.append(error_msg)

                            continue

                   

                    # --- D. TEXT CONTENT SLIDES (Skip - Content Builder will handle) ---

                    else:

                        # Create placeholder slide with title only

                        try:

                            layout_idx = 1  # Title & Content

                            if layout_idx >= len(prs.slide_layouts):

                                layout_idx = min(1, len(prs.slide_layouts) - 1)

                           

                            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

                            slide_mapping[index] = slide

                           

                            # Set title only

                            if slide.shapes.title:

                                slide.shapes.title.text = str(title)

                           

                            slides_created += 1

                            logger.info(f"Created placeholder slide {index + 1}: {title} (content will be added by content builder)")

                           

                        except Exception as e:

                            error_msg = f"Error creating placeholder slide {index + 1}: {e}"

                            logger.error(error_msg, exc_info=True)

                            errors.append(error_msg)

                            continue

                           

                except Exception as e:

                    error_msg = f"Unexpected error processing slide {index + 1}: {e}"

                    logger.error(error_msg, exc_info=True)

                    errors.append(error_msg)

                    continue

                   

        except Exception as e:

            error_msg = f"Critical error during visualization slide generation: {e}"

            logger.error(error_msg, exc_info=True)

            return {

                "temp_pptx_path": "",

                "current_step": "visualization_build_failed",

                "error_message": error_msg

            }

       

        # Validate minimum slide requirement

        if slides_created < 4:

            error_msg = f"Only {slides_created} slides created. Minimum 4 required."

            logger.error(error_msg)

            return {

                "temp_pptx_path": "",

                "current_step": "visualization_build_failed",

                "error_message": error_msg

            }

       

        # Save to temporary file

        temp_path = "temp_presentation.pptx"

        # Apply Branding to ALL slides just before saving
        try:
            total_slides = len(prs.slides)
            for i, slide in enumerate(prs.slides):
                # i+1 for 1-based indexing
                add_branding_to_slide(slide, i + 1, total_slides)
            logger.info("Applied branding (Logo & Page Numbers) to all slides.")
        except Exception as e:
            logger.warning(f"Error applying branding: {e}")

        try:

            prs.save(temp_path)

            logger.info(f"Visualization builder: Saved temporary presentation to {temp_path}")

        except Exception as e:

            error_msg = f"Failed to save temporary presentation: {e}"

            logger.error(error_msg)

            return {

                "temp_pptx_path": "",

                "current_step": "visualization_build_failed",

                "error_message": error_msg

            }

       

        # Verify file was created

        if not os.path.exists(temp_path):

            error_msg = "Temporary presentation file was not created"

            logger.error(error_msg)

            return {

                "temp_pptx_path": "",

                "current_step": "visualization_build_failed",

                "error_message": error_msg

            }

       

        if errors:

            logger.warning(f"Visualization builder: Created presentation with {len(errors)} warnings: {errors}")

       

        logger.info(f"=== VISUALIZATION BUILDER AGENT: Successfully created {slides_created} slides with charts/tables ===")

       

        return {

            "temp_pptx_path": temp_path,

            "current_step": "visualization_build_complete",

            "error_message": "; ".join(errors) if errors else ""

        }

       

    except Exception as e:

        error_msg = f"Visualization builder agent unexpected error: {str(e)}"

        logger.error(error_msg, exc_info=True)

        return {

            "temp_pptx_path": "",

            "current_step": "visualization_build_failed",

            "error_message": error_msg

        }





# --- AGENT 2: CONTENT BUILDER NODE ---

def content_builder_node(state):

    """

    Content Builder Agent: Adds text content to PowerPoint slides and finalizes presentation.

   

    RESPONSIBILITIES:

    - Load temporary PPTX file from visualization builder

    - Add text content to slides that need it

    - Format text content professionally

    - Add speaker notes if available

    - Finalize and save the presentation

   

    INSTRUCTIONS:

    1. Load the temporary PPTX file created by visualization builder

    2. For each slide that has 'content' field (text bullets):

       - Find the appropriate text placeholder

       - Clear existing placeholder content

       - Add each bullet point as a paragraph

       - Format text: size 18pt, proper spacing

       - Ensure professional appearance

    3. For slides with both visualizations and content:

       - Add content below the chart/table if space allows

       - Or add content in a separate text box positioned appropriately

    4. Add speaker notes if provided in slide_data

    5. Ensure all slides are properly formatted

    6. Save final presentation

    7. Clean up temporary file

    """

    logger.info("=== CONTENT BUILDER AGENT: Adding text content and finalizing ===")

   

    try:

        temp_path = state.get("temp_pptx_path", "")

        slides_data = state.get("slides", [])

       

        # Validation

        if not temp_path:

            error_msg = "No temporary PPTX file provided from visualization builder"

            logger.error(error_msg)

            return {

                "final_file_path": "",

                "current_step": "content_build_failed",

                "error_message": error_msg

            }

       

        if not os.path.exists(temp_path):

            error_msg = f"Temporary PPTX file not found: {temp_path}"

            logger.error(error_msg)

            return {

                "final_file_path": "",

                "current_step": "content_build_failed",

                "error_message": error_msg

            }

       

        if not slides_data:

            error_msg = "No slides data provided to content builder"

            logger.error(error_msg)

            return {

                "final_file_path": "",

                "current_step": "content_build_failed",

                "error_message": error_msg

            }

       

        logger.info(f"Content builder: Loading presentation from {temp_path}")

       

        # Load presentation

        try:

            prs = Presentation(temp_path)

            logger.info(f"Successfully loaded presentation with {len(prs.slides)} slides")

        except Exception as e:

            error_msg = f"Failed to load temporary presentation: {e}"

            logger.error(error_msg)

            return {

                "final_file_path": "",

                "current_step": "content_build_failed",

                "error_message": error_msg

            }

       

        # Ensure we have enough slides

        if len(prs.slides) < len(slides_data):

            logger.warning(f"Presentation has {len(prs.slides)} slides but {len(slides_data)} slide data entries")

       

        errors = []

        content_added_count = 0

       

        try:

            for index, slide_data in enumerate(slides_data):

                try:

                    # Calculate actual PPT slide index
                    # Index 0 is Title -> ppt_index 0
                    # Document Control is at ppt_index 1
                    # Index 1+ in slides_data -> ppt_index = index + 1
                    
                    if index == 0:
                        ppt_index = 0
                    else:
                        ppt_index = index + 1

                    if ppt_index >= len(prs.slides):

                        logger.warning(f"Slide index {ppt_index} exceeds available slides. Skipping.")

                        continue

                   

                    if not isinstance(slide_data, dict):

                        continue

                   

                    slide = prs.slides[ppt_index]

                    title = slide_data.get('title', f'Slide {index + 1}')

                    content = slide_data.get('content', [])

                    is_chart = 'chart_type' in slide_data and slide_data.get('chart_type')

                    is_table = 'table_data' in slide_data and slide_data.get('table_data')

                   

                    # Skip title slide (index 0) - subtitle already handled by visualization builder

                    if index == 0:

                        continue

                   

                    # Add text content to slides

                    if content and (not is_chart and not is_table):

                        # Text-only slide - add content to placeholder

                        try:

                            body_shape = None

                            for shape in slide.placeholders:

                                if shape.has_text_frame and shape != slide.shapes.title:

                                    body_shape = shape

                                    break

                            
                            # FALLBACK: If no body placeholder found, create a manual text box
                            if not body_shape:
                                logger.info(f"Slide {index + 1}: No body placeholder found. Creating fallback text box.")
                                left = Inches(1.0)
                                top = Inches(1.5)
                                width = Inches(8.0)
                                height = Inches(5.0)
                                body_shape = slide.shapes.add_textbox(left, top, width, height)

                           

                            if body_shape:

                                tf = body_shape.text_frame

                                tf.clear()

                               

                                # Normalize content to list AND split newlines
                                content_list = []
                                if isinstance(content, str):
                                    # Split by newlines to handle multi-line strings as separate bullets
                                    lines = content.split('\n')
                                    content_list = [line.strip() for line in lines if line.strip()]
                                elif isinstance(content, list):
                                    for item in content:
                                        if isinstance(item, str):
                                            lines = item.split('\n')
                                            content_list.extend([line.strip() for line in lines if line.strip()])
                                        else:
                                            # Convert non-string items to string
                                            content_list.append(str(item))
                                else:
                                    content_list = [str(content)]
                               
                                if content_list:

                                    for point in content_list:

                                        if point:  # Skip empty points

                                            p = tf.add_paragraph()

                                            p.text = str(point)

                                            p.level = 0

                                            p.font.size = Pt(18)

                                            p.space_after = Pt(12)

                                   

                                    content_added_count += 1

                                    logger.info(f"Added text content to slide {index + 1}: {title}")

                                else:

                                    logger.warning(f"Slide {index + 1} ('{title}'): No content provided")

                            else:

                                logger.warning(f"Slide {index + 1} ('{title}'): No body placeholder found for text content")

                               

                        except Exception as e:

                            error_msg = f"Error adding text content to slide {index + 1}: {e}"

                            logger.error(error_msg)

                            errors.append(error_msg)

                   

                    # Add speaker notes if available

                    try:

                        notes_text = slide_data.get('speaker_notes', "")

                        if notes_text and slide.notes_slide and slide.notes_slide.notes_text_frame:

                            slide.notes_slide.notes_text_frame.text = str(notes_text)

                            logger.debug(f"Added speaker notes to slide {index + 1}")

                    except Exception as e:

                        logger.warning(f"Error adding speaker notes to slide {index + 1}: {e}")

                       

                except Exception as e:

                    error_msg = f"Unexpected error processing slide {index + 1}: {e}"

                    logger.error(error_msg, exc_info=True)

                    errors.append(error_msg)

                    continue

                   

        except Exception as e:

            error_msg = f"Critical error during content addition: {e}"

            logger.error(error_msg, exc_info=True)

            return {

                "final_file_path": "",

                "current_step": "content_build_failed",

                "error_message": error_msg

            }

       

        # Validate minimum slide requirement

        if len(prs.slides) < 4:

            error_msg = f"Only {len(prs.slides)} slides in presentation. Minimum 4 required."

            logger.error(error_msg)

            return {

                "final_file_path": "",

                "current_step": "content_build_failed",

                "error_message": error_msg

            }

       

        # Save final presentation

        output_path = "Presentation.pptx"

        try:

            prs.save(output_path)

            logger.info(f"Content builder: Saved final presentation to {output_path}")

        except Exception as e:

            error_msg = f"Failed to save final presentation: {e}"

            logger.error(error_msg)

            return {

                "final_file_path": "",

                "current_step": "content_build_failed",

                "error_message": error_msg

            }

       

        # Verify file was created

        if not os.path.exists(output_path):

            error_msg = "Final presentation file was not created"

            logger.error(error_msg)

            return {

                "final_file_path": "",

                "current_step": "content_build_failed",

                "error_message": error_msg

            }

       

        # Clean up temporary file

        try:

            if os.path.exists(temp_path):

                os.remove(temp_path)

                logger.info(f"Cleaned up temporary file: {temp_path}")

        except Exception as e:

            logger.warning(f"Failed to clean up temporary file: {e}")

       

        if errors:

            logger.warning(f"Content builder: Finalized presentation with {len(errors)} warnings: {errors}")

       

        logger.info(f"=== CONTENT BUILDER AGENT: Successfully finalized presentation with {len(prs.slides)} slides (added content to {content_added_count} slides) ===")

       

        return {

            "final_file_path": output_path,

            "current_step": "content_build_complete",

            "error_message": "; ".join(errors) if errors else ""

        }

       

    except Exception as e:

        error_msg = f"Content builder agent unexpected error: {str(e)}"

        logger.error(error_msg, exc_info=True)

        return {

            "final_file_path": "",

            "current_step": "content_build_failed",

            "error_message": error_msg

        }